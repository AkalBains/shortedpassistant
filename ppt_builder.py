from __future__ import annotations
"""ppt_builder.py – final Python 3.9 version
===========================================
*   Preserves the template’s paragraph‑level font & size for every textbox
*   Moves green rating spheres along the grey track named **track_rating**
*   Resizes energy bars relative to their *original* width (rating 5)
*   Enlarges radar charts to 4×4 in if placeholders are tiny

Rename your shapes once in the template and never touch the code again.
"""

from pathlib import Path
from typing import Dict, List, Union, Optional

from pptx import Presentation
from pptx.util import Inches

# ---------------------------------------------------------------------------
# Helpers – lookup, text, bars, balls
# ---------------------------------------------------------------------------


def _safe_shape(slide, target: str):
    """Return the first shape whose .name matches *target* or *None* and warn."""
    for shp in slide.shapes:
        if shp.name == target:
            return shp
    print(f"[ppt_builder] WARNING – shape '{target}' not found on slide {slide.slide_number}")
    return None


def _set_text_preserve_style(shape, new_text: str) -> None:
    """Overwrite text while copying the font of each original paragraph."""
    if shape is None or not shape.has_text_frame:
        return

    # Snapshot font styles per paragraph (first run per paragraph)
    templates = []
    for p in shape.text_frame.paragraphs:
        if p.runs:
            templates.append(p.runs[0].font)
    if not templates:  # empty textbox
        shape.text = new_text
        return

    shape.text_frame.clear()
    for idx, line in enumerate(new_text.split("\n")):
        p = shape.text_frame.add_paragraph()
        r = p.add_run()
        r.text = line
        tpl = templates[min(idx, len(templates) - 1)]
        r.font.name = tpl.name
        r.font.size = tpl.size
        r.font.bold = tpl.bold
        r.font.italic = tpl.italic
        p.level = 0


# Cache of original full‑widths per bar rectangle
_BAR_ORIG_WIDTH: Dict[str, int] = {}


def _resize_bar_relative(shape, rating: float) -> None:
    """Resize *shape* to rating/5 × original width."""
    if shape is None or shape.has_text_frame:
        return
    if shape.name not in _BAR_ORIG_WIDTH:
        _BAR_ORIG_WIDTH[shape.name] = shape.width
    shape.width = int(_BAR_ORIG_WIDTH[shape.name] * (rating / 5.0))


def _move_ball(ball, rating: int, left_emu: int, right_emu: int) -> None:
    if ball is None:
        return
    frac = (rating - 1) / 4.0
    ball.left = int(left_emu + frac * (right_emu - left_emu))


# ---------------------------------------------------------------------------
# Main builder
# ---------------------------------------------------------------------------

def build_report_pptx(
    template_path: Union[str, Path],
    output_path: Union[str, Path],
    *,
    candidate_name: str,
    role_and_company: str,
    personal_profile: str,
    strengths: List[Dict[str, str]],
    development_areas: List[Dict[str, str]],
    future_considerations: str,
    personal_development: List[Dict[str, str]],
    org_support: List[Dict[str, str]],
    radar_chart_1_path: Optional[Union[str, Path]] = None,
    radar_chart_2_path: Optional[Union[str, Path]] = None,
    bar_scores: Optional[Dict[str, float]] = None,
    summary_ratings: Optional[Dict[str, int]] = None,
    reasoning_scores: Optional[Dict[str, int]] = None,
) -> Path:
    prs = Presentation(str(template_path))

    # --------------------  SLIDE 1  --------------------
    s1 = prs.slides[0]
    _set_text_preserve_style(_safe_shape(s1, "candidate_name"), candidate_name)
    _set_text_preserve_style(_safe_shape(s1, "role_company"), role_and_company)

    # --------------------  SLIDE 3  --------------------
    s3 = prs.slides[2]
    _set_text_preserve_style(_safe_shape(s3, "personal_profile"), personal_profile)

    for i in range(3):
        _set_text_preserve_style(_safe_shape(s3, f"strength_{i+1}_title"), strengths[i]["title"])
        _set_text_preserve_style(_safe_shape(s3, f"strength_{i+1}_body"), strengths[i]["paragraph"])
        _set_text_preserve_style(_safe_shape(s3, f"dev_{i+1}_title"), development_areas[i]["title"])
        _set_text_preserve_style(_safe_shape(s3, f"dev_{i+1}_body"), development_areas[i]["paragraph"])

    _set_text_preserve_style(_safe_shape(s3, "future_considerations"), future_considerations)

    # green balls
    track = _safe_shape(s3, "track_rating")  # grey track rectangle
    if summary_ratings and track:
        left_emu = track.left
        right_emu = track.left + track.width
        _move_ball(_safe_shape(s3, "ball_fit"), summary_ratings["Fit for Role"], left_emu, right_emu)
        _move_ball(_safe_shape(s3, "ball_cap"), summary_ratings["Capabilities"], left_emu, right_emu)
        _move_ball(_safe_shape(s3, "ball_pot"), summary_ratings["Potential"], left_emu, right_emu)
        _move_ball(_safe_shape(s3, "ball_future"), summary_ratings["Future Considerations"], left_emu, right_emu)

    # --------------------  SLIDE 4  --------------------
    s4 = prs.slides[3]
    for i in range(2):
        _set_text_preserve_style(_safe_shape(s4, f"pd_{i+1}_title"), personal_development[i]["title"])
        _set_text_preserve_style(_safe_shape(s4, f"pd_{i+1}_body"), personal_development[i]["paragraph"])
        _set_text_preserve_style(_safe_shape(s4, f"org_{i+1}_title"), org_support[i]["title"])
        _set_text_preserve_style(_safe_shape(s4, f"org_{i+1}_body"), org_support[i]["paragraph"])

    # --------------------  SLIDE 5  --------------------
    if bar_scores:
        s5 = prs.slides[4]
        name_map = {
            "purpose energy": "bar_purpose",
            "intellectual energy": "bar_intellectual",
            "emotional energy": "bar_emotional",
            "people energy": "bar_people",
            "performance impact": "bar_performance",
            "strategic framing": "bar_strategic",
            "mobilisation": "bar_mobilisation",
            "powerful relationships": "bar_relationships",
        }
        for label, score in bar_scores.items():
            shp = _safe_shape(s5, name_map[label.lower()])
            _resize_bar_relative(shp, score)

    # --------------------  SLIDE 6  --------------------
    s6 = prs.slides[5]
    RAD_W = Inches(4.5)
    RAD_H = Inches(4.5)

    def _insert_radar(ph_name: str, png_path):
        ph = _safe_shape(s6, ph_name)
        if ph and Path(png_path).exists():
            w = ph.width if ph.height > Inches(3) else RAD_W
            h = ph.height if ph.height > Inches(3) else RAD_H
            s6.shapes.add_picture(str(png_path), ph.left, ph.top, w, h)
            ph._element.getparent().remove(ph._element)

    if radar_chart_1_path:
        _insert_radar("spider_1", radar_chart_1_path)
    if radar_chart_2_path:
        _insert_radar("spider_2", radar_chart_2_path)

    # --------------------  SLIDE 7  --------------------
    if reasoning_scores:
        s7 = prs.slides[6]
        for label, score in reasoning_scores.items():
            _resize_bar_relative(_safe_shape(s7, f"bar_{label.lower()}"), float(score) / 20)  # 100 → 5
            _set_text_preserve_style(_safe_shape(s7, f"label_{label.lower()}"), f"{score}%")

    # --------------------  SAVE  --------------------
    out_path = Path(output_path)
    out
