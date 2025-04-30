from __future__ import annotations
"""ppt_builder.py
--------------------------------------------------
Fills the cleaned PowerPoint template with narrative, ratings, and images.
All shapes are addressed **by name** (set once in the template’s Selection
Pane).  If a shape is missing the builder skips it rather than crashing.

Usage
-----
Call `build_report_pptx()` from your Streamlit app (see app.py).

Template naming convention
--------------------------
Slide 1
    candidate_name          – text box
    role_company            – text box

Slide 3
    personal_profile        – text box
    strength_1_title        – text box
    strength_1_body         – text box
    (… through 3)
    dev_1_title, dev_1_body – text boxes (… through 3)
    future_considerations   – text box
    ball_fit                – green circle (ellipse)
    ball_cap                – green circle
    ball_pot                – green circle
    ball_future             – green circle
    (you may optionally add underlying grey rectangles named track_fit, etc.)

Slide 4
    pd_1_title, pd_1_body, pd_2_title, pd_2_body
    org_1_title, org_1_body, org_2_title, org_2_body

Slide 5
    bar_purpose, bar_intellectual, bar_emotional, bar_people,
    bar_performance, bar_strategic, bar_mobilisation, bar_relationships

Slide 6
    spider_1, spider_2   – placeholders (text boxes or picture placeholders)

Slide 7
    bar_verbal, bar_numerical, bar_abstract, bar_overall
    label_verbal, label_numerical, label_abstract, label_overall

If any of those names differ, adjust the constants in this file.
"""

from pathlib import Path
from typing import Dict, List, Any
import os

from pptx import Presentation
from pptx.util import Inches

# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------


def _set_text(shape, new_text: str):
    """Safely replace the text in a shape if it exists."""
    if shape is None or not shape.has_text_frame:
        return
    shape.text = new_text


def _safe_shape(slide, name: str):
    """Return shape by name or None (no KeyError)."""
    try:
        return slide.shapes[name]
    except KeyError:
        return None


def _resize_bar(shape, score: float, max_score: float, max_width_emu):
    """Resize a rectangle horizontally according to score (left edge fixed)."""
    if shape is None or shape.has_text_frame:
        return
    shape.width = int(max_width_emu * (score / max_score))


def _reposition_ball(ball, rating: int, track_left_emu: int, track_right_emu: int):
    """Move the ball along a horizontal track based on a 1‒5 rating."""
    if ball is None:
        return
    frac = (rating - 1) / 4  # 0.0 – 1.0
    ball.left = track_left_emu + int((track_right_emu - track_left_emu) * frac)


# ---------------------------------------------------------------------------
# Main public helper
# ---------------------------------------------------------------------------

def build_report_pptx(
    template_path: str | Path,
    output_path: str | Path,
    *,
    candidate_name: str,
    role_and_company: str,
    personal_profile: str,
    strengths: List[Dict[str, str]],
    development_areas: List[Dict[str, str]],
    future_considerations: str,
    personal_development: List[Dict[str, str]],
    org_support: List[Dict[str, str]],
    radar_chart_1_path: str | Path | None = None,
    radar_chart_2_path: str | Path | None = None,
    bar_scores: Dict[str, float] | None = None,
    summary_ratings: Dict[str, int] | None = None,
    reasoning_scores: Dict[str, int] | None = None,
) -> Path:
    """Fill the PowerPoint template and save it.

    Returns the *Path* to the saved file.
    """

    prs = Presentation(str(template_path))

    # ------------------------------------------------------------------
    # Slide 1 – basics
    # ------------------------------------------------------------------
    s1 = prs.slides[0]
    _set_text(_safe_shape(s1, "candidate_name"), candidate_name)
    _set_text(_safe_shape(s1, "role_company"), role_and_company)

    # ------------------------------------------------------------------
    # Slide 3 – narrative & green balls
    # ------------------------------------------------------------------
    s3 = prs.slides[2]
    _set_text(_safe_shape(s3, "personal_profile"), personal_profile)

    for i in range(3):
        _set_text(_safe_shape(s3, f"strength_{i+1}_title"), strengths[i]["title"])
        _set_text(_safe_shape(s3, f"strength_{i+1}_body"), strengths[i]["paragraph"])
        _set_text(_safe_shape(s3, f"dev_{i+1}_title"), development_areas[i]["title"])
        _set_text(_safe_shape(s3, f"dev_{i+1}_body"), development_areas[i]["paragraph"])

    _set_text(_safe_shape(s3, "future_considerations"), future_considerations)

    # green balls – horizontal track positions (adjust if your template differs)
    if summary_ratings:
        TRACK_LEFT = Inches(1.2).emu
        TRACK_RIGHT = Inches(6.5).emu
        _reposition_ball(_safe_shape(s3, "ball_fit"), summary_ratings.get("Fit for Role", 3), TRACK_LEFT, TRACK_RIGHT)
        _reposition_ball(_safe_shape(s3, "ball_cap"), summary_ratings.get("Capabilities", 3), TRACK_LEFT, TRACK_RIGHT)
        _reposition_ball(_safe_shape(s3, "ball_pot"), summary_ratings.get("Potential", 3), TRACK_LEFT, TRACK_RIGHT)
        _reposition_ball(_safe_shape(s3, "ball_future"), summary_ratings.get("Future Considerations", 3), TRACK_LEFT, TRACK_RIGHT)

    # ------------------------------------------------------------------
    # Slide 4 – development suggestions
    # ------------------------------------------------------------------
    s4 = prs.slides[3]
    for i in range(2):
        _set_text(_safe_shape(s4, f"pd_{i+1}_title"), personal_development[i]["title"])
        _set_text(_safe_shape(s4, f"pd_{i+1}_body"), personal_development[i]["paragraph"])
        _set_text(_safe_shape(s4, f"org_{i+1}_title"), org_support[i]["title"])
        _set_text(_safe_shape(s4, f"org_{i+1}_body"), org_support[i]["paragraph"])

    # ------------------------------------------------------------------
    # Slide 5 – eight energy bars
    # ------------------------------------------------------------------
    if bar_scores:
        s5 = prs.slides[4]
        max_width = Inches(5.5).emu
        name_map = {
            "purpose energy": "bar_purpose",
            "intellectual energy": "bar_intellectual",
            "emotional energy": "bar_emotional",
            "people energy": "bar_people",
            "performance impact": "bar_performance",
            "strategic framing": "bar_strategic",
            "mobilisation": "bar_mobilisation",
            "powerful relationships": "bar_relationships",
        }
        for label, score in bar_scores.items():
            shape_name = name_map.get(label.lower())
            if not shape_name:
                continue
            _resize_bar(_safe_shape(s5, shape_name), score, 5, max_width)

    # ------------------------------------------------------------------
    # Slide 6 – radar charts
    # ------------------------------------------------------------------
    s6 = prs.slides[5]
    if radar_chart_1_path and Path(radar_chart_1_path).exists():
        ph = _safe_shape(s6, "spider_1")
        if ph:
            s6.shapes.add_picture(str(radar_chart_1_path), ph.left, ph.top, ph.width, ph.height)
            # remove placeholder
            ph._element.getparent().remove(ph._element)
    if radar_chart_2_path and Path(radar_chart_2_path).exists():
        ph = _safe_shape(s6, "spider_2")
        if ph:
            s6.shapes.add_picture(str(radar_chart_2_path), ph.left, ph.top, ph.width, ph.height)
            ph._element.getparent().remove(ph._element)

    # ------------------------------------------------------------------
    # Slide 7 – reasoning bars
    # ------------------------------------------------------------------
    if reasoning_scores:
        s7 = prs.slides[6]
        max_width = Inches(5.5).emu
        for label, score in reasoning_scores.items():
            name = label.lower()
            _resize_bar(_safe_shape(s7, f"bar_{name}"), score, 100, max_width)
            _set_text(_safe_shape(s7, f"label_{name}"), f"{score}%")

    # ------------------------------------------------------------------
    # Save
    # ------------------------------------------------------------------
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path
